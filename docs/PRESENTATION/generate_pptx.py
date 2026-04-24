import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

C_PRIMARY = RGBColor(0xC4, 0x95, 0x6A)
C_PRIMARY_DARK = RGBColor(0x9E, 0x76, 0x4E)
C_DARK = RGBColor(0x1A, 0x1A, 0x2E)
C_DARK2 = RGBColor(0x16, 0x21, 0x3E)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT = RGBColor(0xF5, 0xF0, 0xEB)
C_GRAY = RGBColor(0x99, 0x99, 0x99)
C_LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
C_ACCENT = RGBColor(0xF6, 0xAD, 0x55)
C_GREEN = RGBColor(0x48, 0xBB, 0x78)
C_RED = RGBColor(0xE5, 0x3E, 0x3E)
C_BLUE = RGBColor(0x41, 0x82, 0xC4)

FONT_CN = "微软雅黑"
FONT_EN = "Calibri"
FONT_NUM = "DIN"


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None, shape_type=MSO_SHAPE.RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text="", font_size=18, font_color=C_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT_CN, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_paragraph(text_frame, text, font_size=18, font_color=C_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT_CN, space_before=0, space_after=4, line_spacing=1.2):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return p


def add_accent_line(slide, left, top, width, color=C_PRIMARY, height=Pt(4)):
    line = add_shape(slide, left, top, width, height, fill_color=color)
    return line


def add_page_number(slide, num, total=15):
    add_textbox(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
                f"{num}/{total}", font_size=10, font_color=C_GRAY, alignment=PP_ALIGN.RIGHT, font_name=FONT_EN)


def add_top_bar(slide):
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), fill_color=C_PRIMARY)


def add_bottom_bar(slide):
    add_shape(slide, Inches(0), Inches(7.2), SLIDE_WIDTH, Inches(0.3), fill_color=C_DARK)
    add_textbox(slide, Inches(0.5), Inches(7.2), Inches(4), Inches(0.3),
                "寻裳 XUNO — AI穿搭决策平台", font_size=9, font_color=C_GRAY, font_name=FONT_CN)


def create_card(slide, left, top, width, height, fill_color=None, border_color=None):
    shape = add_shape(slide, left, top, width, height, fill_color=fill_color or RGBColor(0x22, 0x22, 0x3A), line_color=border_color or C_PRIMARY_DARK, line_width=Pt(1))
    shape.shadow.inherit = False
    return shape


def add_big_number(slide, left, top, number, label, num_color=C_PRIMARY, label_color=C_LIGHT_GRAY, num_size=56, label_size=14):
    add_textbox(slide, left, top, Inches(3), Inches(0.9), number, font_size=num_size, font_color=num_color, bold=True, font_name=FONT_NUM, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, top + Inches(0.85), Inches(3), Inches(0.5), label, font_size=label_size, font_color=label_color, alignment=PP_ALIGN.CENTER)


def add_icon_bullet(slide, left, top, icon_text, bullet_text, icon_color=C_PRIMARY, text_color=C_WHITE, text_size=16):
    add_textbox(slide, left, top, Inches(0.4), Inches(0.4), icon_text, font_size=16, font_color=icon_color, alignment=PP_ALIGN.CENTER, font_name=FONT_EN)
    add_textbox(slide, left + Inches(0.5), top, Inches(5), Inches(0.4), bullet_text, font_size=text_size, font_color=text_color)


prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

blank_layout = prs.slide_layouts[6]

# ==================== Page 1: 封面 ====================
slide1 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide1, C_DARK)

add_shape(slide1, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, fill_color=C_DARK2)

add_shape(slide1, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT, fill_color=C_PRIMARY)

add_shape(slide1, Inches(0.5), Inches(2.0), Inches(0.8), Inches(0.8), fill_color=C_PRIMARY, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_textbox(slide1, Inches(0.55), Inches(2.05), Inches(0.7), Inches(0.7), "X", font_size=36, font_color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name=FONT_EN)

add_textbox(slide1, Inches(1.6), Inches(1.8), Inches(8), Inches(1.2), "寻裳 XUNO", font_size=52, font_color=C_WHITE, bold=True)
add_textbox(slide1, Inches(1.6), Inches(2.9), Inches(8), Inches(0.6), "AI穿搭决策平台", font_size=28, font_color=C_PRIMARY)

add_accent_line(slide1, Inches(1.6), Inches(3.7), Inches(2), C_PRIMARY)

add_textbox(slide1, Inches(1.6), Inches(4.0), Inches(8), Inches(0.5), "让每个年轻人都能自信地迈出职场第一步", font_size=20, font_color=C_LIGHT_GRAY)

add_textbox(slide1, Inches(1.6), Inches(5.2), Inches(6), Inches(0.4), "1人 + AI辅助开发  |  21天  |  34万行代码  |  345次提交", font_size=14, font_color=C_GRAY)

add_shape(slide1, Inches(9.5), Inches(1.5), Inches(3.2), Inches(3.2), fill_color=RGBColor(0x1E, 0x1E, 0x38), line_color=C_PRIMARY_DARK, line_width=Pt(2), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_textbox(slide1, Inches(9.5), Inches(2.2), Inches(3.2), Inches(1.5), "[ AI生成面试场景\n  穿搭效果图 ]", font_size=16, font_color=C_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide1, Inches(9.5), Inches(5.0), Inches(3.2), Inches(0.4), "[ 扫码体验Demo ]", font_size=13, font_color=C_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide1, Inches(0.5), Inches(6.8), Inches(6), Inches(0.4), "第十一届中国国际大学生创新大赛  校赛", font_size=12, font_color=C_GRAY)

add_page_number(slide1, 1)

# ==================== Page 2: 痛点 — 小张的故事 ====================
slide2 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide2, C_DARK)
add_top_bar(slide2)
add_bottom_bar(slide2)

add_textbox(slide2, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7), "痛点：小张的故事", font_size=32, font_color=C_WHITE, bold=True)
add_accent_line(slide2, Inches(0.8), Inches(1.1), Inches(1.5), C_PRIMARY)

story_box = create_card(slide2, Inches(0.8), Inches(1.5), Inches(5.5), Inches(4.8))
txBox = add_textbox(slide2, Inches(1.2), Inches(1.8), Inches(4.8), Inches(4.2), "", font_size=17, font_color=C_WHITE, line_spacing=1.5)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
items = [
    ("小张，23岁，后天去互联网公司面试产品经理", C_WHITE),
    ("", C_WHITE),
    ("打开衣柜，30件衣服，不知道穿哪件", C_LIGHT_GRAY),
    ("", C_WHITE),
    ("网上搜\"面试穿什么\" → 2万个结果，越看越慌", C_LIGHT_GRAY),
    ("", C_WHITE),
    ("最后选了最安全的方案 —— 西装", C_ACCENT),
    ("", C_WHITE),
    ("但互联网公司不需要西装", C_RED),
]
for i, (text, color) in enumerate(items):
    if i == 0:
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(17)
        tf.paragraphs[0].font.color.rgb = color
        tf.paragraphs[0].font.name = FONT_CN
        tf.paragraphs[0].space_after = Pt(2)
    else:
        add_paragraph(tf, text, font_size=17, font_color=color, space_after=2, line_spacing=1.5)

create_card(slide2, Inches(6.8), Inches(1.5), Inches(5.5), Inches(2.2), fill_color=RGBColor(0x1E, 0x1E, 0x38))
add_textbox(slide2, Inches(7.2), Inches(1.8), Inches(4.8), Inches(0.4), "[ 衣柜照片：衣服很多但无从选择 ]", font_size=13, font_color=C_GRAY, alignment=PP_ALIGN.CENTER)

create_card(slide2, Inches(6.8), Inches(4.0), Inches(5.5), Inches(2.3), fill_color=RGBColor(0x1E, 0x1E, 0x38))
add_textbox(slide2, Inches(7.2), Inches(4.3), Inches(4.8), Inches(0.4), "[ 搜索截图：\"面试穿什么\" 2万个结果 ]", font_size=13, font_color=C_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide2, Inches(0.8), Inches(6.5), Inches(10), Inches(0.4), "这不是小张一个人的问题，这是900万毕业生的共同困境", font_size=16, font_color=C_ACCENT, bold=True)

add_page_number(slide2, 2)

# ==================== Page 3: 问题量化 ====================
slide3 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide3, C_DARK)
add_top_bar(slide3)
add_bottom_bar(slide3)

add_textbox(slide3, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7), "问题量化：900万人的穿搭焦虑", font_size=32, font_color=C_WHITE, bold=True)
add_accent_line(slide3, Inches(0.8), Inches(1.1), Inches(1.5), C_PRIMARY)

nums = [
    ("900万", "毕业生/年", C_RED),
    ("90%", "不知道面试穿什么", C_ACCENT),
    ("2小时", "平均决定穿搭时间", C_PRIMARY),
    ("1620万", "小时/年 穿搭焦虑", C_RED),
]
for i, (num, label, color) in enumerate(nums):
    left = Inches(0.8 + i * 3.1)
    create_card(slide3, left, Inches(1.6), Inches(2.8), Inches(2.5), border_color=color)
    add_big_number(slide3, left, Inches(1.9), num, label, num_color=color, num_size=48, label_size=14)

add_textbox(slide3, Inches(0.8), Inches(4.5), Inches(11), Inches(0.5), "男性用户尤甚：现有穿搭内容严重偏向女性，男性几乎被忽视", font_size=16, font_color=C_LIGHT_GRAY)

create_card(slide3, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.2), fill_color=RGBColor(0x1E, 0x1E, 0x38))
txBox3 = add_textbox(slide3, Inches(1.2), Inches(5.4), Inches(10.8), Inches(0.8), "", font_size=12, font_color=C_GRAY)
tf3 = txBox3.text_frame
tf3.paragraphs[0].text = "数据来源："
tf3.paragraphs[0].font.size = Pt(12)
tf3.paragraphs[0].font.color.rgb = C_GRAY
tf3.paragraphs[0].font.name = FONT_CN
add_paragraph(tf3, "• 900万毕业生 — 教育部2025数据", font_size=12, font_color=C_GRAY, space_before=2)
add_paragraph(tf3, "• 90%不知道面试穿什么 — 智联招聘《2025中国大学生就业报告》", font_size=12, font_color=C_GRAY, space_before=2)

add_page_number(slide3, 3)

# ==================== Page 4: 解决方案 — 伊伊AI造型师 ====================
slide4 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide4, C_DARK)
add_top_bar(slide4)
add_bottom_bar(slide4)

add_textbox(slide4, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "解决方案：伊伊AI造型师", font_size=32, font_color=C_WHITE, bold=True)
add_accent_line(slide4, Inches(0.8), Inches(1.1), Inches(1.5), C_PRIMARY)

add_textbox(slide4, Inches(0.8), Inches(1.4), Inches(6), Inches(0.5), "首页即入口，无需搜索 — 3步解决穿搭焦虑", font_size=18, font_color=C_PRIMARY)

steps = [
    ("01", "告诉伊伊场景", "\"明天去互联网公司面试\"", C_PRIMARY),
    ("02", "获取个性化推荐", "3套Smart Casual搭配\n+ 匹配度分析", C_ACCENT),
    ("03", "AI试穿确认", "虚拟试穿效果\n确认后一键购买", C_GREEN),
]
for i, (num, title, desc, color) in enumerate(steps):
    top = Inches(2.1 + i * 1.5)
    add_shape(slide4, Inches(0.8), top, Inches(0.7), Inches(0.7), fill_color=color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide4, Inches(0.8), top + Inches(0.05), Inches(0.7), Inches(0.6), num, font_size=22, font_color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name=FONT_EN)
    add_textbox(slide4, Inches(1.7), top + Inches(0.02), Inches(3.5), Inches(0.4), title, font_size=18, font_color=C_WHITE, bold=True)
    add_textbox(slide4, Inches(1.7), top + Inches(0.4), Inches(3.5), Inches(0.6), desc, font_size=13, font_color=C_LIGHT_GRAY)

    if i < 2:
        add_shape(slide4, Inches(1.1), top + Inches(0.75), Inches(0.05), Inches(0.7), fill_color=C_PRIMARY_DARK)

create_card(slide4, Inches(6.5), Inches(1.6), Inches(5.8), Inches(3.8), fill_color=RGBColor(0x1E, 0x1E, 0x38))
add_textbox(slide4, Inches(6.5), Inches(2.5), Inches(5.8), Inches(2.0), "[ 产品界面截图\n伊伊对话界面\n面试场景推荐流程 ]", font_size=16, font_color=C_GRAY, alignment=PP_ALIGN.CENTER)

create_card(slide4, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.0), fill_color=RGBColor(0x1E, 0x1E, 0x38), border_color=C_GREEN)
add_textbox(slide4, Inches(1.2), Inches(5.9), Inches(4), Inches(0.4), "Before: 搜索2小时 + 可能穿错", font_size=15, font_color=C_RED)
add_textbox(slide4, Inches(6.5), Inches(5.9), Inches(0.5), Inches(0.4), "→", font_size=20, font_color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide4, Inches(7.2), Inches(5.9), Inches(4), Inches(0.4), "After: XUNO 3分钟 + 穿得对", font_size=15, font_color=C_GREEN, bold=True)

add_page_number(slide4, 4)

# ==================== Page 5: 产品Demo ====================
slide5 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide5, C_DARK)
add_top_bar(slide5)
add_bottom_bar(slide5)

add_textbox(slide5, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7), "产品Demo：核心体验", font_size=32, font_color=C_WHITE, bold=True)
add_accent_line(slide5, Inches(0.8), Inches(1.1), Inches(1.5), C_PRIMARY)

add_textbox(slide5, Inches(0.8), Inches(1.4), Inches(6), Inches(0.5), "2分20秒 Demo精华 — 3个\"哇\"时刻", font_size=18, font_color=C_PRIMARY)

create_card(slide5, Inches(0.8), Inches(2.1), Inches(11.5), Inches(4.0), fill_color=RGBColor(0x0D, 0x0D, 0x1A), border_color=C_PRIMARY_DARK)
add_textbox(slide5, Inches(0.8), Inches(3.2), Inches(11.5), Inches(1.5), "[ Demo 视频播放区域 ]\n\n手机投屏 / 实时演示 / 录屏播放", font_size=20, font_color=C_GRAY, alignment=PP_ALIGN.CENTER)

wow_items = [
    ("哇1", "语音交互", "伊伊主动问候，无需搜索", C_PRIMARY),
    ("哇2", "匹配度雷达图", "五维可视化 + 实时优化", C_ACCENT),
    ("哇3", "记忆能力", "跨场景偏好记忆", C_GREEN),
]
for i, (label, title, desc, color) in enumerate(wow_items):
    left = Inches(0.8 + i * 4.0)
    create_card(slide5, left, Inches(6.3), Inches(3.6), Inches(0.7), border_color=color)
    add_textbox(slide5, left + Inches(0.2), Inches(6.35), Inches(0.7), Inches(0.3), label, font_size=12, font_color=color, bold=True)
    add_textbox(slide5, left + Inches(0.9), Inches(6.35), Inches(2.5), Inches(0.3), title, font_size=14, font_color=C_WHITE, bold=True)
    add_textbox(slide5, left + Inches(0.9), Inches(6.6), Inches(2.5), Inches(0.3), desc, font_size=11, font_color=C_LIGHT_GRAY)

add_page_number(slide5, 5)

# ==================== Page 6: 技术架构 ====================
slide6 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide6, C_DARK)
add_top_bar(slide6)
add_bottom_bar(slide6)

add_textbox(slide6, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), "技术架构：FashionCLIP + 对话状态机 + 六层推荐", font_size=30, font_color=C_WHITE, bold=True)
add_accent_line(slide6, Inches(0.8), Inches(1.1), Inches(1.5), C_PRIMARY)

funnel_layers = [
    ("规则引擎", "场景/体型/肤色基础规则", C_PRIMARY),
    ("场景匹配", "FashionCLIP 512维向量语义匹配", C_PRIMARY_DARK),
    ("体型适配", "5种体型分类 + 阈值配置", RGBColor(0xB8, 0x86, 0x4A)),
    ("色彩协调", "肤色-服装色彩协调度计算", C_ACCENT),
    ("风格一致性", "用户风格偏好向量匹配", RGBColor(0xD4, 0xA5, 0x74)),
    ("LLM精排", "GLM-4-Flash 最终排序+解释", C_GREEN),
]
for i, (name, desc, color) in enumerate(funnel_layers):
    top = Inches(1.5 + i * 0.85)
    w = Inches(5.5 - i * 0.3)
    left = Inches(0.8 + i * 0.15)
    add_shape(slide6, left, top, w, Inches(0.7), fill_color=color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide6, left + Inches(0.2), top + Inches(0.05), Inches(2), Inches(0.3), name, font_size=14, font_color=C_WHITE, bold=True)
    add_textbox(slide6, left + Inches(0.2), top + Inches(0.35), w - Inches(0.4), Inches(0.3), desc, font_size=11, font_color=C_WHITE)

create_card(slide6, Inches(7.0), Inches(1.5), Inches(5.5), Inches(2.3), fill_color=RGBColor(0x1E, 0x1E, 0x38))
add_textbox(slide6, Inches(7.3), Inches(1.6), Inches(5), Inches(0.3), "FashionCLIP 视觉理解", font_size=14, font_color=C_PRIMARY, bold=True)
add_textbox(slide6, Inches(7.3), Inches(2.0), Inches(5), Inches(1.5), "服装图片 → 512维向量 → 语义匹配\n理解\"oversized blazer\" vs \"slim-fit shirt\"\n比通用CLIP准确率高~18%", font_size=12, font_color=C_LIGHT_GRAY, line_spacing=1.5)

create_card(slide6, Inches(7.0), Inches(4.0), Inches(5.5), Inches(1.5), fill_color=RGBColor(0x1E, 0x1E, 0x38))
add_textbox(slide6, Inches(7.3), Inches(4.1), Inches(5), Inches(0.3), "对话状态机", font_size=14, font_color=C_PRIMARY, bold=True)
add_textbox(slide6, Inches(7.3), Inches(4.5), Inches(5), Inches(0.8), "GREET(问候) → CONTEXT(采集) → GENERATE(推荐)\nRedis存储 | TTL 3600秒 | 即时状态更新", font_size=12, font_color=C_LIGHT_GRAY, line_spacing=1.5)

create_card(slide6, Inches(7.0), Inches(5.7), Inches(5.5), Inches(1.2), fill_color=RGBColor(0x1E, 0x1E, 0x38))
add_textbox(slide6, Inches(7.3), Inches(5.8), Inches(5), Inches(0.3), "多LLM降级链", font_size=14, font_color=C_PRIMARY, bold=True)
add_textbox(slide6, Inches(7.3), Inches(6.2), Inches(5), Inches(0.6), "GLM(主) → DeepSeek(备) → 豆包(备)\n熔断器: 5次失败→熔断 | 60秒后半开", font_size=12, font_color=C_LIGHT_GRAY, line_spacing=1.5)

add_page_number(slide6, 6)

# ==================== Page 7: 创新点1 — ChineseFashionCLIP ====================
slide7 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide7, C_DARK)
add_top_bar(slide7)
add_bottom_bar(slide7)

add_textbox(slide7, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), "创新点1：ChineseFashionCLIP 微调", font_size=32, font_color=C_WHITE, bold=True)
add_accent_line(slide7, Inches(0.8), Inches(1.1), Inches(1.5), C_PRIMARY)

add_textbox(slide7, Inches(0.8), Inches(1.4), Inches(10), Inches(0.5), "问题：FashionCLIP基于Farfetch(英国)数据训练，有西方审美偏见", font_size=16, font_color=C_LIGHT_GRAY)

create_card(slide7, Inches(0.8), Inches(2.1), Inches(5.5), Inches(2.5), fill_color=RGBColor(0x2A, 0x1A, 0x1A), border_color=C_RED)
add_textbox(slide7, Inches(1.2), Inches(2.2), Inches(4.8), Inches(0.3), "原版 FashionCLIP", font_size=14, font_color=C_RED, bold=True)
add_textbox(slide7, Inches(1.2), Inches(2.6), Inches(4.8), Inches(0.3), "输入: \"学院风面试搭配\"", font_size=13, font_color=C_LIGHT_GRAY)
add_textbox(slide7, Inches(1.2), Inches(3.0), Inches(4.8), Inches(0.3), "返回: 西装革履（西方语境）", font_size=15, font_color=C_WHITE, bold=True)
add_textbox(slide7, Inches(1.2), Inches(3.5), Inches(4.8), Inches(0.6), "[ 截图位置：原版检索结果 ]", font_size=12, font_color=C_GRAY, alignment=PP_ALIGN.CENTER)

create_card(slide7, Inches(6.8), Inches(2.1), Inches(5.5), Inches(2.5), fill_color=RGBColor(0x1A, 0x2A, 0x1A), border_color=C_GREEN)
add_textbox(slide7, Inches(7.2), Inches(2.2), Inches(4.8), Inches(0.3), "ChineseFashionCLIP 微调版", font_size=14, font_color=C_GREEN, bold=True)
add_textbox(slide7, Inches(7.2), Inches(2.6), Inches(4.8), Inches(0.3), "输入: \"学院风面试搭配\"", font_size=13, font_color=C_LIGHT_GRAY)
add_textbox(slide7, Inches(7.2), Inches(3.0), Inches(4.8), Inches(0.3), "返回: 休闲西装+白T+牛仔裤", font_size=15, font_color=C_WHITE, bold=True)
add_textbox(slide7, Inches(7.2), Inches(3.5), Inches(4.8), Inches(0.6), "[ 截图位置：微调版检索结果 ]", font_size=12, font_color=C_GRAY, alignment=PP_ALIGN.CENTER)

add_big_number(slide7, Inches(0.8), Inches(5.0), "~18%", "检索准确率提升", num_color=C_GREEN, num_size=48, label_size=14)
add_textbox(slide7, Inches(3.8), Inches(5.2), Inches(4), Inches(0.4), "来源: EVIDENCE/decisions.md ADR-001", font_size=11, font_color=C_GRAY)

create_card(slide7, Inches(6.8), Inches(4.8), Inches(5.5), Inches(2.0), fill_color=RGBColor(0x1E, 0x1E, 0x38))
add_textbox(slide7, Inches(7.2), Inches(4.9), Inches(4.8), Inches(0.3), "数据飞轮", font_size=14, font_color=C_PRIMARY, bold=True)
add_textbox(slide7, Inches(7.2), Inches(5.3), Inches(4.8), Inches(1.2), "用户使用 → 反馈数据 → 模型优化 → 推荐更准 → 用户更依赖\n\n含金量：这是我们自己训练的模型，不是API调用", font_size=13, font_color=C_LIGHT_GRAY, line_spacing=1.4)

add_page_number(slide7, 7)

# ==================== Page 8: 创新点2 — 对话式记忆 ====================
slide8 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide8, C_DARK)
add_top_bar(slide8)
add_bottom_bar(slide8)

add_textbox(slide8, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "创新点2：对话式记忆", font_size=32, font_color=C_WHITE, bold=True)
add_accent_line(slide8, Inches(0.8), Inches(1.1), Inches(1.5), C_PRIMARY)

add_textbox(slide8, Inches(0.8), Inches(1.4), Inches(8), Inches(0.5), "伊伊记住你在对话中表达的每一个偏好 — 越用越懂你", font_size=18, font_color=C_PRIMARY)

create_card(slide8, Inches(0.8), Inches(2.1), Inches(6.0), Inches(3.5), fill_color=RGBColor(0x1E, 0x1E, 0x38))
add_textbox(slide8, Inches(1.2), Inches(2.2), Inches(5.2), Inches(0.3), "场景演示", font_size=14, font_color=C_PRIMARY, bold=True)

chat_items = [
    ("用户:", "我不喜欢高领的", C_WHITE),
    ("伊伊:", "记住了！以后都推荐圆领和V领给你 ✓", C_GREEN),
    ("", "", C_WHITE),
    ("[ 切换到约会场景 ]", "", C_ACCENT),
    ("伊伊:", "推荐V领针织衫 — 因为你说过不喜欢高领", C_GREEN),
]
for i, (speaker, text, color) in enumerate(chat_items):
    top = Inches(2.7 + i * 0.45)
    if speaker:
        add_textbox(slide8, Inches(1.4), top, Inches(0.8), Inches(0.35), speaker, font_size=13, font_color=C_LIGHT_GRAY, bold=True)
        add_textbox(slide8, Inches(2.2), top, Inches(4.2), Inches(0.35), text, font_size=13, font_color=color)
    else:
        add_textbox(slide8, Inches(1.4), top, Inches(4.8), Inches(0.35), text, font_size=13, font_color=color)

create_card(slide8, Inches(7.3), Inches(2.1), Inches(5.2), Inches(3.5), fill_color=RGBColor(0x1E, 0x1E, 0x38))
add_textbox(slide8, Inches(7.7), Inches(2.2), Inches(4.5), Inches(0.3), "技术实现", font_size=14, font_color=C_PRIMARY, bold=True)
tech_items = [
    "• 跨场景应用: \"不喜欢高领\"在面试/约会/日常都生效",
    "• Redis存储: dialog:context:{sessionId}",
    "• TTL 3600秒 + 即时状态更新",
    "• 不是静态画像，是实时对话理解",
    "• 关键词提取 → 偏好存储 → 推荐过滤",
]
for i, item in enumerate(tech_items):
    add_textbox(slide8, Inches(7.7), Inches(2.7 + i * 0.5), Inches(4.5), Inches(0.45), item, font_size=12, font_color=C_LIGHT_GRAY, line_spacing=1.3)

add_textbox(slide8, Inches(0.8), Inches(6.0), Inches(11), Inches(0.5), "核心壁垒：记忆不是预设标签，是实时对话理解 — 竞争对手做不到", font_size=16, font_color=C_ACCENT, bold=True)

add_page_number(slide8, 8)

# ==================== Page 9: 创新点3 — 穿搭偏好进化 ====================
slide9 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide9, C_DARK)
add_top_bar(slide9)
add_bottom_bar(slide9)

add_textbox(slide9, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "创新点3：穿搭偏好进化", font_size=32, font_color=C_WHITE, bold=True)
add_accent_line(slide9, Inches(0.8), Inches(1.1), Inches(1.5), C_PRIMARY)

add_textbox(slide9, Inches(0.8), Inches(1.4), Inches(8), Inches(0.5), "用户的穿搭偏好不是静态的，而是持续进化的", font_size=18, font_color=C_PRIMARY)

radar_stages = [
    ("Day 1", "偏好不确定", "雷达图较小\n探索阶段", RGBColor(0x66, 0x66, 0x88)),
    ("Day 7", "偏好开始成形", "方向明确\n风格初现", C_PRIMARY_DARK),
    ("Day 30", "偏好清晰饱满", "伊伊比你更\n了解你的风格", C_PRIMARY),
]
for i, (day, status, desc, color) in enumerate(radar_stages):
    left = Inches(0.8 + i * 4.0)
    create_card(slide9, left, Inches(2.2), Inches(3.6), Inches(2.8), border_color=color)
    add_textbox(slide9, left, Inches(2.4), Inches(3.6), Inches(0.4), day, font_size=24, font_color=color, bold=True, alignment=PP_ALIGN.CENTER, font_name=FONT_EN)
    add_textbox(slide9, left, Inches(2.9), Inches(3.6), Inches(0.4), status, font_size=16, font_color=C_WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide9, left, Inches(3.5), Inches(3.6), Inches(1.0), "[ 风格偏好雷达图 ]\n" + desc, font_size=12, font_color=C_GRAY, alignment=PP_ALIGN.CENTER)

    if i < 2:
        add_textbox(slide9, Inches(4.4 + i * 4.0), Inches(3.2), Inches(0.5), Inches(0.4), "→", font_size=24, font_color=C_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

create_card(slide9, Inches(0.8), Inches(5.3), Inches(11.5), Inches(1.5), fill_color=RGBColor(0x1E, 0x1E, 0x38), border_color=C_PRIMARY)
add_textbox(slide9, Inches(1.2), Inches(5.4), Inches(4), Inches(0.3), "数据飞轮", font_size=14, font_color=C_PRIMARY, bold=True)
add_textbox(slide9, Inches(1.2), Inches(5.8), Inches(10.5), Inches(0.8), "用户使用 → 数据积累 → 模型优化 → 推荐更准 → 用户更依赖 → 30天后，伊伊比你更了解你的风格", font_size=14, font_color=C_LIGHT_GRAY, line_spacing=1.4)

add_page_number(slide9, 9)

# ==================== Page 10: 用户数据 ====================
slide10 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide10, C_DARK)
add_top_bar(slide10)
add_bottom_bar(slide10)

add_textbox(slide10, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "用户数据：种子用户测试结果", font_size=32, font_color=C_WHITE, bold=True)
add_accent_line(slide10, Inches(0.8), Inches(1.1), Inches(1.5), C_PRIMARY)

metrics = [
    ("SUS", "XX/100", "系统可用性评分", C_PRIMARY),
    ("满意度", "XX%", "推荐满意度", C_ACCENT),
    ("决策时间", "X→Y分钟", "平均决策时间缩短", C_GREEN),
    ("NPS", "XX", "净推荐值", C_BLUE),
]
for i, (label, value, desc, color) in enumerate(metrics):
    left = Inches(0.8 + i * 3.1)
    create_card(slide10, left, Inches(1.6), Inches(2.8), Inches(2.2), border_color=color)
    add_textbox(slide10, left, Inches(1.8), Inches(2.8), Inches(0.3), label, font_size=14, font_color=C_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide10, left, Inches(2.2), Inches(2.8), Inches(0.7), value, font_size=40, font_color=color, bold=True, alignment=PP_ALIGN.CENTER, font_name=FONT_NUM)
    add_textbox(slide10, left, Inches(3.0), Inches(2.8), Inches(0.3), desc, font_size=12, font_color=C_GRAY, alignment=PP_ALIGN.CENTER)

create_card(slide10, Inches(0.8), Inches(4.2), Inches(7.0), Inches(2.5), fill_color=RGBColor(0x1E, 0x1E, 0x38))
add_textbox(slide10, Inches(1.2), Inches(4.3), Inches(6.2), Inches(0.3), "典型用户反馈", font_size=14, font_color=C_PRIMARY, bold=True)
feedbacks = [
    "\"以前面试穿搭要纠结一晚上，用了XUNO三分钟就搞定了\" — 小张",
    "\"终于有个App不只会推荐女装\" — 小李",
    "\"伊伊真的记得我说过的话，太神奇了\" — 小王",
]
for i, fb in enumerate(feedbacks):
    add_textbox(slide10, Inches(1.2), Inches(4.8 + i * 0.55), Inches(6.2), Inches(0.5), fb, font_size=13, font_color=C_LIGHT_GRAY, line_spacing=1.3)

create_card(slide10, Inches(8.2), Inches(4.2), Inches(4.1), Inches(2.5), fill_color=RGBColor(0x1E, 0x1E, 0x38))
add_textbox(slide10, Inches(8.2), Inches(5.0), Inches(4.1), Inches(1.0), "[ 用户使用App\n真实照片 ]", font_size=14, font_color=C_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide10, Inches(0.8), Inches(6.9), Inches(10), Inches(0.3), "数据来源于5人种子用户测试，2026年5月 | 数据待填充: prompts/15-user-test.md", font_size=11, font_color=C_GRAY)

add_page_number(slide10, 10)

# ==================== Page 11: 商业模式 ====================
slide11 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide11, C_DARK)
add_top_bar(slide11)
add_bottom_bar(slide11)

add_textbox(slide11, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "商业模式：精益创业，低盈亏平衡点", font_size=32, font_color=C_WHITE, bold=True)
add_accent_line(slide11, Inches(0.8), Inches(1.1), Inches(1.5), C_PRIMARY)

create_card(slide11, Inches(0.8), Inches(1.5), Inches(5.5), Inches(4.0), fill_color=RGBColor(0x1E, 0x1E, 0x38))
add_textbox(slide11, Inches(1.2), Inches(1.6), Inches(4.8), Inches(0.3), "收入模型", font_size=16, font_color=C_PRIMARY, bold=True)

revenue_items = [
    ("免费层", "AI推荐 + 基础对话", C_GREEN, "0元"),
    ("内容产品", "色彩报告6.9元 | 体型分析9.9元\n胶囊衣橱14.9元", C_ACCENT, "付费"),
    ("电商佣金", "淘宝客2-5% + 京东8-15%\n混合加权约5-8%", C_PRIMARY, "交易"),
]
for i, (title, desc, color, price) in enumerate(revenue_items):
    top = Inches(2.1 + i * 1.1)
    add_shape(slide11, Inches(1.2), top, Inches(0.1), Inches(0.8), fill_color=color)
    add_textbox(slide11, Inches(1.5), top, Inches(2), Inches(0.3), title, font_size=14, font_color=color, bold=True)
    add_textbox(slide11, Inches(1.5), top + Inches(0.3), Inches(3.5), Inches(0.5), desc, font_size=12, font_color=C_LIGHT_GRAY, line_spacing=1.3)
    add_textbox(slide11, Inches(4.8), top + Inches(0.1), Inches(1.2), Inches(0.3), price, font_size=13, font_color=color, bold=True, alignment=PP_ALIGN.RIGHT)

create_card(slide11, Inches(6.8), Inches(1.5), Inches(5.5), Inches(4.0), fill_color=RGBColor(0x1E, 0x1E, 0x38))
add_textbox(slide11, Inches(7.2), Inches(1.6), Inches(4.8), Inches(0.3), "成本结构（1000 DAU/月）", font_size=16, font_color=C_PRIMARY, bold=True)

cost_items = [
    ("基础设施", "350元", "服务器+数据库+存储"),
    ("虚拟试穿API", "1,125元", "GLM-4V-Plus调用"),
    ("LLM备用", "100元", "DeepSeek/Qwen备用"),
    ("杂项", "50元", "域名+监控+其他"),
]
for i, (name, cost, desc) in enumerate(cost_items):
    top = Inches(2.2 + i * 0.6)
    add_textbox(slide11, Inches(7.2), top, Inches(2), Inches(0.3), name, font_size=13, font_color=C_LIGHT_GRAY)
    add_textbox(slide11, Inches(9.5), top, Inches(1.2), Inches(0.3), cost, font_size=13, font_color=C_WHITE, bold=True, alignment=PP_ALIGN.RIGHT)
    add_textbox(slide11, Inches(10.8), top, Inches(1.5), Inches(0.3), desc, font_size=10, font_color=C_GRAY)

add_shape(slide11, Inches(7.2), Inches(4.7), Inches(4.8), Pt(1), fill_color=C_PRIMARY_DARK)
add_textbox(slide11, Inches(7.2), Inches(4.8), Inches(2), Inches(0.3), "月总成本", font_size=14, font_color=C_WHITE, bold=True)
add_textbox(slide11, Inches(9.5), Inches(4.8), Inches(1.2), Inches(0.3), "~1,625元", font_size=14, font_color=C_PRIMARY, bold=True, alignment=PP_ALIGN.RIGHT)

highlight_nums = [
    ("月成本", "~1,625元", C_PRIMARY),
    ("盈亏平衡", "5,000 MAU", C_GREEN),
    ("LLM成本", "0元", C_GREEN),
]
for i, (label, value, color) in enumerate(highlight_nums):
    left = Inches(0.8 + i * 4.2)
    create_card(slide11, left, Inches(5.8), Inches(3.8), Inches(1.0), border_color=color)
    add_textbox(slide11, left, Inches(5.9), Inches(3.8), Inches(0.3), label, font_size=12, font_color=C_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide11, left, Inches(6.2), Inches(3.8), Inches(0.4), value, font_size=24, font_color=color, bold=True, alignment=PP_ALIGN.CENTER, font_name=FONT_NUM)

add_page_number(slide11, 11)

# ==================== Page 12: 市场规模 ====================
slide12 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide12, C_DARK)
add_top_bar(slide12)
add_bottom_bar(slide12)

add_textbox(slide12, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "市场规模：TAM / SAM / SOM", font_size=32, font_color=C_WHITE, bold=True)
add_accent_line(slide12, Inches(0.8), Inches(1.1), Inches(1.5), C_PRIMARY)

add_shape(slide12, Inches(2.5), Inches(1.6), Inches(5.0), Inches(5.0), fill_color=RGBColor(0x1A, 0x1A, 0x30), line_color=C_GRAY, line_width=Pt(2), shape_type=MSO_SHAPE.OVAL)
add_textbox(slide12, Inches(2.5), Inches(2.8), Inches(5.0), Inches(0.5), "TAM", font_size=16, font_color=C_GRAY, bold=True, alignment=PP_ALIGN.CENTER, font_name=FONT_EN)
add_textbox(slide12, Inches(2.5), Inches(3.3), Inches(5.0), Inches(0.5), "2.3万亿/年", font_size=22, font_color=C_GRAY, bold=True, alignment=PP_ALIGN.CENTER, font_name=FONT_NUM)
add_textbox(slide12, Inches(2.5), Inches(3.8), Inches(5.0), Inches(0.4), "中国服装电商市场", font_size=12, font_color=C_GRAY, alignment=PP_ALIGN.CENTER)

add_shape(slide12, Inches(3.2), Inches(2.2), Inches(3.6), Inches(3.6), fill_color=RGBColor(0x22, 0x22, 0x40), line_color=C_PRIMARY_DARK, line_width=Pt(2), shape_type=MSO_SHAPE.OVAL)
add_textbox(slide12, Inches(3.2), Inches(3.1), Inches(3.6), Inches(0.4), "SAM", font_size=14, font_color=C_PRIMARY_DARK, bold=True, alignment=PP_ALIGN.CENTER, font_name=FONT_EN)
add_textbox(slide12, Inches(3.2), Inches(3.5), Inches(3.6), Inches(0.4), "4,600亿/年", font_size=18, font_color=C_PRIMARY_DARK, bold=True, alignment=PP_ALIGN.CENTER, font_name=FONT_NUM)
add_textbox(slide12, Inches(3.2), Inches(3.9), Inches(3.6), Inches(0.3), "18-35岁线上消费者", font_size=10, font_color=C_PRIMARY_DARK, alignment=PP_ALIGN.CENTER)

add_shape(slide12, Inches(3.8), Inches(2.7), Inches(2.4), Inches(2.4), fill_color=RGBColor(0x33, 0x33, 0x55), line_color=C_PRIMARY, line_width=Pt(2), shape_type=MSO_SHAPE.OVAL)
add_textbox(slide12, Inches(3.8), Inches(3.3), Inches(2.4), Inches(0.3), "SOM", font_size=12, font_color=C_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER, font_name=FONT_EN)
add_textbox(slide12, Inches(3.8), Inches(3.6), Inches(2.4), Inches(0.3), "1亿/年", font_size=16, font_color=C_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER, font_name=FONT_NUM)
add_textbox(slide12, Inches(3.8), Inches(3.9), Inches(2.4), Inches(0.3), "50万用户×200元ARPU", font_size=9, font_color=C_PRIMARY, alignment=PP_ALIGN.CENTER)

milestones = [
    ("2026 Q3", "1,000用户", "校赛→省赛验证"),
    ("2027 Q1", "10,000用户", "种子轮"),
    ("2027 Q4", "50,000用户", "A轮"),
    ("2028", "500,000用户", "规模化"),
]
add_textbox(slide12, Inches(8.5), Inches(1.6), Inches(4), Inches(0.3), "增长路径", font_size=16, font_color=C_PRIMARY, bold=True)
for i, (time, users, stage) in enumerate(milestones):
    top = Inches(2.2 + i * 1.0)
    add_shape(slide12, Inches(8.5), top + Inches(0.15), Inches(0.15), Inches(0.15), fill_color=C_PRIMARY, shape_type=MSO_SHAPE.OVAL)
    add_textbox(slide12, Inches(8.9), top, Inches(1.5), Inches(0.3), time, font_size=13, font_color=C_ACCENT, bold=True, font_name=FONT_EN)
    add_textbox(slide12, Inches(10.5), top, Inches(1.5), Inches(0.3), users, font_size=13, font_color=C_WHITE, bold=True)
    add_textbox(slide12, Inches(8.9), top + Inches(0.3), Inches(3.5), Inches(0.3), stage, font_size=11, font_color=C_GRAY)
    if i < 3:
        add_shape(slide12, Inches(8.55), top + Inches(0.35), Inches(0.05), Inches(0.6), fill_color=C_PRIMARY_DARK)

add_textbox(slide12, Inches(8.5), Inches(6.2), Inches(4.5), Inches(0.5), "对标: Stitch Fix（美国）年营收21亿美元\n验证AI穿搭推荐市场需求真实存在", font_size=12, font_color=C_LIGHT_GRAY, line_spacing=1.4)

add_page_number(slide12, 12)

# ==================== Page 13: 竞品分析 ====================
slide13 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide13, C_DARK)
add_top_bar(slide13)
add_bottom_bar(slide13)

add_textbox(slide13, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "竞品分析：独特优势", font_size=32, font_color=C_WHITE, bold=True)
add_accent_line(slide13, Inches(0.8), Inches(1.1), Inches(1.5), C_PRIMARY)

add_textbox(slide13, Inches(0.8), Inches(1.4), Inches(10), Inches(0.5), "本质区别：从\"你去找穿搭\"变成\"穿搭来找你\"", font_size=18, font_color=C_PRIMARY)

headers = ["维度", "淘宝", "小红书", "抖音", "XUNO"]
col_widths = [Inches(1.8), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.8)]
col_colors = [C_LIGHT_GRAY, C_GRAY, C_GRAY, C_GRAY, C_PRIMARY]

table_left = Inches(0.8)
table_top = Inches(2.1)
row_height = Inches(0.65)

for j, (header, w, color) in enumerate(zip(headers, col_widths, col_colors)):
    left = table_left + sum(cw for cw in [0] + [col_widths[k] for k in range(j)])
    add_shape(slide13, int(left), table_top, int(w), int(row_height), fill_color=RGBColor(0x2A, 0x2A, 0x45) if j < 4 else C_PRIMARY_DARK)
    add_textbox(slide13, int(left) + Inches(0.1), table_top + Inches(0.1), int(w) - Inches(0.2), int(row_height) - Inches(0.2), header, font_size=14, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)

rows_data = [
    ["首页入口", "4层深埋", "搜索触发", "无", "首页即入口"],
    ["推荐方式", "商品列表", "UGC内容", "短视频", "AI对话"],
    ["记忆能力", "无", "无", "无", "穿搭记忆"],
    ["试穿体验", "静态图片", "无", "无", "AI生成"],
    ["包容性", "分男女", "偏女性", "偏女性", "全人群"],
]
for i, row in enumerate(rows_data):
    y = table_top + row_height * (i + 1)
    bg_color = RGBColor(0x1E, 0x1E, 0x38) if i % 2 == 0 else RGBColor(0x22, 0x22, 0x3E)
    for j, (cell, w, color) in enumerate(zip(row, col_widths, [C_LIGHT_GRAY, C_WHITE, C_WHITE, C_WHITE, C_PRIMARY])):
        left = table_left + sum(cw for cw in [0] + [col_widths[k] for k in range(j)])
        cell_bg = bg_color if j < 4 else RGBColor(0x1A, 0x2A, 0x1A)
        add_shape(slide13, int(left), int(y), int(w), int(row_height), fill_color=cell_bg)
        is_bold = j == 0 or j == 4
        add_textbox(slide13, int(left) + Inches(0.1), int(y) + Inches(0.1), int(w) - Inches(0.2), int(row_height) - Inches(0.2), cell, font_size=13, font_color=color, bold=is_bold, alignment=PP_ALIGN.CENTER)

add_textbox(slide13, Inches(0.8), Inches(5.8), Inches(11), Inches(0.5), "XUNO独有: 对话式推荐 + 穿搭记忆 + AI试穿 + 全人群覆盖", font_size=16, font_color=C_PRIMARY, bold=True)

create_card(slide13, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.6), fill_color=RGBColor(0x1E, 0x1E, 0x38), border_color=C_PRIMARY)
add_textbox(slide13, Inches(1.2), Inches(6.45), Inches(10.5), Inches(0.4), "\"只有XUNO让AI主动找到你，而不是你去搜\"", font_size=15, font_color=C_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

add_page_number(slide13, 13)

# ==================== Page 14: 团队 + 路线图 ====================
slide14 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide14, C_DARK)
add_top_bar(slide14)
add_bottom_bar(slide14)

add_textbox(slide14, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "团队 + 路线图", font_size=32, font_color=C_WHITE, bold=True)
add_accent_line(slide14, Inches(0.8), Inches(1.1), Inches(1.5), C_PRIMARY)

create_card(slide14, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.2), fill_color=RGBColor(0x1E, 0x1E, 0x38))
add_textbox(slide14, Inches(1.2), Inches(1.6), Inches(5), Inches(0.3), "团队", font_size=16, font_color=C_PRIMARY, bold=True)

add_textbox(slide14, Inches(1.2), Inches(2.1), Inches(5), Inches(0.3), "袁荣跃 — 独立全栈开发者", font_size=15, font_color=C_WHITE, bold=True)
add_textbox(slide14, Inches(1.2), Inches(2.5), Inches(5), Inches(0.3), "开发方式: AI辅助开发（Claude Code + GLM）", font_size=12, font_color=C_LIGHT_GRAY)

dev_stats = [
    ("21", "天"),
    ("345", "次提交"),
    ("34万", "行代码"),
    ("1,581", "个源文件"),
]
for i, (num, label) in enumerate(dev_stats):
    left = Inches(1.2 + i * 1.3)
    add_textbox(slide14, left, Inches(3.0), Inches(1.2), Inches(0.5), num, font_size=28, font_color=C_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER, font_name=FONT_NUM)
    add_textbox(slide14, left, Inches(3.5), Inches(1.2), Inches(0.3), label, font_size=11, font_color=C_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide14, Inches(1.2), Inches(4.1), Inches(5), Inches(0.3), "AI辅助下开发效率约5.5倍", font_size=14, font_color=C_GREEN, bold=True)
add_textbox(slide14, Inches(1.2), Inches(4.5), Inches(5), Inches(0.3), "5个关键决策ADR记录 | AI写代码，我做决策", font_size=12, font_color=C_LIGHT_GRAY)

ai_contrib = [
    ("代码生成", "60%", "AI / 40% 人"),
    ("架构设计", "30%", "AI / 70% 人"),
    ("文档", "80%", "AI / 20% 人"),
    ("决策", "10%", "AI / 90% 人"),
]
for i, (area, ai_pct, split) in enumerate(ai_contrib):
    top = Inches(5.0 + i * 0.4)
    add_textbox(slide14, Inches(1.2), top, Inches(1.5), Inches(0.3), area, font_size=11, font_color=C_LIGHT_GRAY)
    bar_w = float(ai_pct.replace("%", "")) / 100.0 * 2.5
    add_shape(slide14, Inches(2.8), top + Inches(0.05), Inches(bar_w), Inches(0.2), fill_color=C_PRIMARY)
    add_shape(slide14, Inches(2.8 + bar_w), top + Inches(0.05), Inches(2.5 - bar_w), Inches(0.2), fill_color=RGBColor(0x33, 0x33, 0x55))
    add_textbox(slide14, Inches(5.5), top, Inches(1.5), Inches(0.3), split, font_size=10, font_color=C_GRAY)

create_card(slide14, Inches(7.0), Inches(1.5), Inches(5.5), Inches(5.2), fill_color=RGBColor(0x1E, 0x1E, 0x38))
add_textbox(slide14, Inches(7.4), Inches(1.6), Inches(4.8), Inches(0.3), "路线图", font_size=16, font_color=C_PRIMARY, bold=True)

roadmap = [
    ("2026.05", "校赛（现在）", True, C_ACCENT),
    ("2026.07", "省赛 + 种子用户100人", False, C_PRIMARY),
    ("2026.09", "正式上线", False, C_PRIMARY),
    ("2026.12", "1,000用户", False, C_PRIMARY),
    ("2027.Q1", "种子轮融资", False, C_GREEN),
]
for i, (time, milestone, is_current, color) in enumerate(roadmap):
    top = Inches(2.2 + i * 0.85)
    dot_size = Inches(0.2) if is_current else Inches(0.12)
    dot_left = Inches(7.6) if is_current else Inches(7.64)
    add_shape(slide14, dot_left, top + Inches(0.08), dot_size, dot_size, fill_color=color, shape_type=MSO_SHAPE.OVAL)
    if i < 4:
        add_shape(slide14, Inches(7.67), top + Inches(0.3), Inches(0.05), Inches(0.55), fill_color=C_PRIMARY_DARK)
    add_textbox(slide14, Inches(8.0), top, Inches(1.5), Inches(0.3), time, font_size=13, font_color=color, bold=True, font_name=FONT_EN)
    add_textbox(slide14, Inches(9.5), top, Inches(2.8), Inches(0.3), milestone, font_size=13, font_color=C_WHITE if is_current else C_LIGHT_GRAY, bold=is_current)
    if is_current:
        add_textbox(slide14, Inches(9.5), top + Inches(0.3), Inches(2.8), Inches(0.2), "← 我们在这里", font_size=10, font_color=C_ACCENT)

add_textbox(slide14, Inches(7.4), Inches(6.2), Inches(4.8), Inches(0.3), "正在寻找: 技术合伙人 + 运营合伙人", font_size=13, font_color=C_ACCENT, bold=True)

add_page_number(slide14, 14)

# ==================== Page 15: Q&A ====================
slide15 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide15, C_DARK2)

add_shape(slide15, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT, fill_color=C_PRIMARY)

add_textbox(slide15, Inches(1.6), Inches(1.2), Inches(8), Inches(0.8), "感谢聆听", font_size=44, font_color=C_WHITE, bold=True)
add_accent_line(slide15, Inches(1.6), Inches(2.1), Inches(2), C_PRIMARY)

txBox15 = add_textbox(slide15, Inches(1.6), Inches(2.5), Inches(8), Inches(1.5), "", font_size=17, font_color=C_LIGHT_GRAY, line_spacing=1.6)
tf15 = txBox15.text_frame
tf15.paragraphs[0].text = "小张用了XUNO，3分钟搞定面试穿搭。"
tf15.paragraphs[0].font.size = Pt(17)
tf15.paragraphs[0].font.color.rgb = C_LIGHT_GRAY
tf15.paragraphs[0].font.name = FONT_CN
add_paragraph(tf15, "面试那天，他自信地走进大楼——", font_size=17, font_color=C_LIGHT_GRAY, space_before=4)
add_paragraph(tf15, "不是因为西装穿得好，而是因为穿得对。", font_size=17, font_color=C_PRIMARY, bold=True, space_before=4)
add_paragraph(tf15, "", font_size=10, space_before=8)
add_paragraph(tf15, "900万个小张，都值得这份自信。", font_size=18, font_color=C_ACCENT, bold=True, space_before=8)

create_card(slide15, Inches(8.5), Inches(1.5), Inches(3.8), Inches(3.8), fill_color=RGBColor(0x1E, 0x1E, 0x38), border_color=C_PRIMARY)
add_textbox(slide15, Inches(8.5), Inches(2.5), Inches(3.8), Inches(2.0), "[ 扫码体验\n你的第一套\n面试穿搭 ]", font_size=16, font_color=C_GRAY, alignment=PP_ALIGN.CENTER)

seek_items = [
    "100名大学生免费体验面试穿搭",
    "导师: 时尚行业/投资领域",
    "合作: 高校就业指导中心",
]
for i, item in enumerate(seek_items):
    add_textbox(slide15, Inches(1.6), Inches(4.8 + i * 0.4), Inches(6), Inches(0.35), f"• {item}", font_size=14, font_color=C_LIGHT_GRAY)

add_textbox(slide15, Inches(1.6), Inches(6.2), Inches(4), Inches(0.3), "袁荣跃 | 寻裳 XUNO", font_size=14, font_color=C_WHITE, bold=True)
add_textbox(slide15, Inches(1.6), Inches(6.5), Inches(4), Inches(0.3), "联系方式: [微信/邮箱]", font_size=12, font_color=C_GRAY)

add_textbox(slide15, Inches(8.5), Inches(6.0), Inches(3.8), Inches(0.5), "欢迎提问", font_size=24, font_color=C_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

add_page_number(slide15, 15)

output_path = r"C:\AiNeed\docs\PRESENTATION\XUNO-FINAL.pptx"
prs.save(output_path)
print(f"PPTX saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
